"""첨부 파일에서 텍스트 추출 (txt, eml, Office 일부)."""

from __future__ import annotations

import email.policy
import io
from email import message_from_bytes


def extract_text_from_bytes(filename: str, data: bytes, max_chars: int = 80_000) -> str:
    name = (filename or "").lower()
    if not data:
        return ""

    if name.endswith(".txt") or name.endswith(".csv") or name.endswith(".log"):
        return _decode_plain(data, max_chars)

    if name.endswith(".eml"):
        return _extract_eml(data, max_chars)

    if name.endswith(".docx"):
        return _extract_docx(data, max_chars)

    if name.endswith(".xlsx") or name.endswith(".xlsm"):
        return _extract_xlsx(data, max_chars)

    if name.endswith(".pptx"):
        return _extract_pptx(data, max_chars)

    if name.endswith(".doc") or name.endswith(".ppt") or name.endswith(".xls"):
        return f"[{filename}] 구형 Office(.doc/.xls/.ppt)는 이 환경에서 텍스트 추출을 지원하지 않습니다. txt/docx/xlsx/pptx/eml을 사용해 주세요."

    return f"[{filename}] 지원 형식: txt, csv, eml, docx, xlsx, pptx."


def _decode_plain(data: bytes, max_chars: int) -> str:
    for enc in ("utf-8", "utf-8-sig", "cp949", "euc-kr", "latin-1"):
        try:
            t = data.decode(enc)
            break
        except UnicodeDecodeError:
            t = ""
    else:
        t = data.decode("utf-8", errors="replace")
    return t[:max_chars]


def _extract_eml(data: bytes, max_chars: int) -> str:
    try:
        msg = message_from_bytes(data, policy=email.policy.default)
    except Exception as exc:  # noqa: BLE001
        return f"[eml 파싱 오류] {exc}"
    parts: list[str] = []
    subj = msg.get("subject")
    if subj:
        parts.append(f"Subject: {subj}")
    frm = msg.get("from")
    if frm:
        parts.append(f"From: {frm}")
    to = msg.get("to")
    if to:
        parts.append(f"To: {to}")
    date = msg.get("date")
    if date:
        parts.append(f"Date: {date}")

    body_chunks: list[str] = []

    def walk(part) -> None:
        ctype = part.get_content_type()
        if ctype == "text/plain":
            try:
                body_chunks.append(part.get_content())
            except Exception:  # noqa: BLE001
                pass
        elif ctype == "text/html" and not body_chunks:
            try:
                raw = part.get_content()
                body_chunks.append(_strip_html_simple(raw))
            except Exception:  # noqa: BLE001
                pass
        elif part.is_multipart():
            for sub in part.iter_parts():
                walk(sub)

    if msg.is_multipart():
        walk(msg)
    else:
        walk(msg)

    parts.append("\n".join(body_chunks))
    out = "\n".join(p for p in parts if p)
    return out[:max_chars]


def _strip_html_simple(html: str) -> str:
    import re

    t = re.sub(r"(?is)<script.*?>.*?</script>", " ", html)
    t = re.sub(r"(?is)<style.*?>.*?</style>", " ", t)
    t = re.sub(r"(?is)<[^>]+>", " ", t)
    return " ".join(t.split())[:50_000]


def _extract_docx(data: bytes, max_chars: int) -> str:
    try:
        from docx import Document
    except ImportError:
        return "[docx] python-docx 패키지가 필요합니다. pip install python-docx"
    try:
        doc = Document(io.BytesIO(data))
        paras = [p.text for p in doc.paragraphs if p.text.strip()]
        return "\n".join(paras)[:max_chars]
    except Exception as exc:  # noqa: BLE001
        return f"[docx 오류] {exc}"


def _extract_xlsx(data: bytes, max_chars: int) -> str:
    try:
        import openpyxl
    except ImportError:
        return "[xlsx] openpyxl 패키지가 필요합니다. pip install openpyxl"
    try:
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        lines: list[str] = []
        for sheet in wb.worksheets:
            lines.append(f"## {sheet.title}")
            for row in sheet.iter_rows(max_row=500, values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(x.strip() for x in cells):
                    lines.append("\t".join(cells))
            lines.append("")
        wb.close()
        return "\n".join(lines)[:max_chars]
    except Exception as exc:  # noqa: BLE001
        return f"[xlsx 오류] {exc}"


def _extract_pptx(data: bytes, max_chars: int) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return "[pptx] python-pptx 패키지가 필요합니다. pip install python-pptx"
    try:
        prs = Presentation(io.BytesIO(data))
        chunks: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text.strip())
            if texts:
                chunks.append(f"--- 슬라이드 {i} ---\n" + "\n".join(texts))
        return "\n\n".join(chunks)[:max_chars]
    except Exception as exc:  # noqa: BLE001
        return f"[pptx 오류] {exc}"
